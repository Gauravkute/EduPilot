from pptx import Presentation

def extract(filepath: str) -> dict:
    prs = Presentation(filepath)
    slides = []

    for idx,slide in enumerate(prs.slides,start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        
        slides.append({"slide": idx, "texts": texts, "notes": notes, "text_count": len(texts)})
        
    return{
        "slide_count":len(prs.slides),
        "slides":slides
    }